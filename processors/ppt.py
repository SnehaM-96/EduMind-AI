from pptx import Presentation


def extract_ppt_text(ppt_file):

    try:

        prs = Presentation(ppt_file)
        text = ""

        for slide_num, slide in enumerate(prs.slides):

            text += f"\n[Slide {slide_num + 1}]\n"

            for shape in slide.shapes:

                try:

                    # Extract text frames
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text.strip() + "\n"

                    # Extract tables inside slides
                    if shape.has_table:

                        for row in shape.table.rows:

                            row_text = " | ".join(
                                cell.text.strip()
                                for cell in row.cells
                                if cell.text.strip()
                            )

                            if row_text:
                                text += row_text + "\n"

                except Exception:
                    continue

        return text.strip()

    except Exception :
        return ""